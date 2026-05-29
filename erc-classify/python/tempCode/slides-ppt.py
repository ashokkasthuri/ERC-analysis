'''
Author: ashokkasthuri ashokk@smu.edu.sg
Date: 2025-07-26 14:18:06
LastEditors: ashokkasthuri ashokk@smu.edu.sg
LastEditTime: 2025-07-26 14:18:12
FilePath: /ERC-analysis-master/erc-classify/slides-ppt.py
Description: 这是默认设置,请设置`customMade`, 打开koroFileHeader查看配置 进行设置: https://github.com/OBKoro1/koro1FileHeader/wiki/%E9%85%8D%E7%BD%AE
'''
from pptx import Presentation
from pptx.util import Pt

title = "Punishment Mechanisms in Web3 — Research Questions"

bullets = [
    ("RQ1 – Functionalities:",
     "What functions do punishment mechanisms serve for trust & security? (Motivation: establish core dimensions & roles)"),
    ("RQ2 – Taxonomy:",
     "Which punishment mechanisms & penalty types are used today in Web3? (Motivation: structured classification)"),
    ("RQ3 – Scenario Mapping:",
     "Which mechanisms are invoked in which application scenarios / malpractices? (Motivation: link misuse → penalties)"),
    ("RQ4 – Enforcement Practice:",
     "How are these mechanisms implemented & enforced (governance, compliance, adaptability)?"),
    ("RQ5 – Open Challenges:",
     "What gaps remain and how far has literature addressed them? (Motivation: future directions)"),
]

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = title
slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

tf = slide.placeholders[1].text_frame
tf.clear()
for head, body in bullets:
    p = tf.add_paragraph()
    p.level = 0
    run = p.add_run()
    run.text = f"{head} {body}"
    run.font.size = Pt(20)

prs.save("Punishment_Mechanisms_Web3_RQs.pptx")
print("Saved: Punishment_Mechanisms_Web3_RQs.pptx")
